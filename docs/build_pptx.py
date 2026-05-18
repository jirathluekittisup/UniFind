"""Build UniFind Proposal Presentation using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
SKY    = RGBColor(0x0e, 0xa5, 0xe9)   # primary sky blue
AMBER  = RGBColor(0xf5, 0x9e, 0x0b)   # accent amber
DARK   = RGBColor(0x0f, 0x17, 0x2a)   # near-black
SLATE  = RGBColor(0x1e, 0x29, 0x3b)   # body text dark
MUTED  = RGBColor(0x64, 0x74, 0x8b)   # muted grey
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BGLT   = RGBColor(0xf0, 0xf9, 0xff)   # very light sky tint
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xe2, 0xe8, 0xf0)
GREEN  = RGBColor(0x16, 0xa3, 0x4a)
RED    = RGBColor(0xdc, 0x26, 0x26)

W, H = Inches(13.33), Inches(7.5)     # LAYOUT_WIDE 16:9

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return txb

def add_textbox_rich(slide, lines, x, y, w, h, default_size=13, default_color=None, font="Calibri"):
    """lines = list of dicts: {text, size, bold, color, italic, align, space_before, bullet}"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        if line.get("space_before"):
            p.space_before = Pt(line["space_before"])
        if line.get("bullet"):
            p.level = line.get("level", 0)
        run = p.add_run()
        run.text = line.get("text", "")
        run.font.size = Pt(line.get("size", default_size))
        run.font.bold = line.get("bold", False)
        run.font.italic = line.get("italic", False)
        run.font.name = font
        col = line.get("color", default_color)
        if col:
            run.font.color.rgb = col
    return txb

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_label_chip(slide, text, x, y, bg_color, text_color=WHITE, size=10):
    """Small pill/chip label."""
    add_rect(slide, x, y, len(text)*0.09+0.25, 0.28, fill=bg_color)
    add_text(slide, text, x+0.06, y+0.03, len(text)*0.09+0.15, 0.22,
             size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – Title
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, DARK)

# Sky-blue left stripe
add_rect(sl, 0, 0, 0.45, 7.5, fill=SKY)

# Amber accent bar
add_rect(sl, 0.45, 2.9, 5.5, 0.07, fill=AMBER)

# App name
add_text(sl, "UniFind", 0.7, 1.5, 7, 1.2, size=72, bold=True, color=WHITE,
         align=PP_ALIGN.LEFT, font="Georgia")

# Subtitle
add_text(sl, "Lost & Found App for University Students",
         0.7, 2.95, 8, 0.55, size=22, bold=False, color=rgb(186,230,253),
         align=PP_ALIGN.LEFT)

# Course + Group
add_text(sl, "2190512 Application Development",
         0.7, 3.7, 7, 0.4, size=14, color=rgb(148,163,184), align=PP_ALIGN.LEFT)
add_text(sl, "[Group Name]  |  April 2026",
         0.7, 4.1, 7, 0.4, size=14, color=rgb(148,163,184), align=PP_ALIGN.LEFT)

# Right decorative circles
add_rect(sl, 10.2, 0.4, 2.5, 2.5, fill=rgb(7,89,133))
add_rect(sl, 11.4, 2.6, 1.8, 1.8, fill=rgb(14,165,233))
add_text(sl, "🔍", 11.6, 3.05, 1.4, 0.9, size=40, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – Agenda
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)

# Top bar
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_text(sl, "Agenda", 0.5, 0.2, 12, 0.65, size=32, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.88, 1.5, 0.04, fill=AMBER)

sections = [
    ("01", "Background & Problems"),
    ("02", "Scope of Application"),
    ("03", "Application Features & Functionalities"),
    ("04", "Limitations of the Application"),
    ("05", "Technology Stack Used"),
]

ystart = 1.2
for i, (num, label) in enumerate(sections):
    y = ystart + i * 1.0
    # Number circle
    add_rect(sl, 0.5, y, 0.75, 0.72, fill=SKY)
    add_text(sl, num, 0.5, y+0.1, 0.75, 0.52, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Label
    add_rect(sl, 1.35, y, 10.5, 0.72, fill=BGLT)
    add_text(sl, label, 1.55, y+0.13, 10.1, 0.48, size=18, bold=False, color=SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – Background & Problems
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

# Section chip
add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=SKY)
add_text(sl, "01", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Background & Problems", 0.5, 0.18, 10, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 2.5, 0.04, fill=AMBER)

problems = [
    "University campuses are large — students frequently lose wallets, keys, student IDs, laptops, AirPods, and more",
    "Current solutions are fragmented: Facebook groups, paper notice boards, security office drop-offs",
    "No unified, searchable, campus-specific platform exists",
    "Finders and losers cannot find each other efficiently",
    "Privacy risks: students must post personal details on public social media groups",
    "Security offices have no digital integration with students",
]

lines = []
for i, p in enumerate(problems):
    lines.append({"text": f"  •  {p}", "size": 13.5, "color": SLATE,
                  "space_before": 4 if i > 0 else 0})

add_textbox_rich(sl, lines, 0.5, 0.9, 8.0, 4.8)

# Pain-point box (right side)
add_rect(sl, 8.9, 0.9, 3.9, 2.5, fill=rgb(254,243,199))
add_rect(sl, 8.9, 0.9, 0.12, 2.5, fill=AMBER)
add_text(sl, "💡 Key Insight", 9.12, 1.0, 3.5, 0.4, size=12, bold=True, color=rgb(146,64,14))
add_textbox_rich(sl, [
    {"text": '"Items get lost forever —', "size": 13, "bold": True, "italic": True, "color": rgb(120,53,15)},
    {"text": "not because no one found them,", "size": 13, "italic": True, "color": rgb(120,53,15)},
    {"text": 'but because there\'s no easy way', "size": 13, "italic": True, "color": rgb(120,53,15)},
    {"text": 'to connect finders and losers."', "size": 13, "bold": True, "italic": True, "color": rgb(120,53,15)},
], 9.12, 1.45, 3.5, 1.8)

# Stats at bottom
stats = [("300+", "Items lost per\nsemester (est.)"), ("< 20%", "Recovery rate\nwithout a system"), ("48h+", "Average wait to\nresolve a case")]
for i, (val, lbl) in enumerate(stats):
    xb = 0.5 + i * 2.7
    add_rect(sl, xb, 5.7, 2.4, 1.55, fill=BGLT)
    add_text(sl, val, xb, 5.82, 2.4, 0.7, size=28, bold=True, color=SKY, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, xb, 6.45, 2.4, 0.7, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – Scope of Application
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=SKY)
add_text(sl, "02", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Scope of Application", 0.5, 0.18, 10, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 2.2, 0.04, fill=AMBER)

# LEFT: In Scope
add_rect(sl, 0.5, 0.95, 5.85, 6.2, fill=rgb(240,249,255))
add_rect(sl, 0.5, 0.95, 5.85, 0.52, fill=SKY)
add_text(sl, "✅  In Scope (v1)", 0.7, 0.98, 5.5, 0.44, size=15, bold=True, color=WHITE)

in_scope = [
    "Single university campus",
    "Verified university students and staff only",
    "iOS and Android mobile app (React Native)",
    "Web admin dashboard for security office",
    "Item reporting (lost & found) with photos",
    "Search and browse listings with filters",
    "In-app messaging between finder and loser",
    "Claim and ownership verification flow",
    "Push notifications for matches and updates",
]
lines_in = []
for i, s in enumerate(in_scope):
    lines_in.append({"text": f"  •  {s}", "size": 13, "color": SLATE, "space_before": 3 if i > 0 else 2})
add_textbox_rich(sl, lines_in, 0.6, 1.55, 5.6, 5.4)

# RIGHT: Out of Scope
add_rect(sl, 6.98, 0.95, 5.85, 6.2, fill=rgb(255,249,245))
add_rect(sl, 6.98, 0.95, 5.85, 0.52, fill=rgb(239,68,68))
add_text(sl, "🚫  Out of Scope (v1)", 7.18, 0.98, 5.5, 0.44, size=15, bold=True, color=WHITE)

out_scope = [
    "Multi-campus or multi-institution support",
    "External / non-university users",
    "Monetary transactions or rewards",
    "AI / computer vision image matching",
    "Real-time WebSocket chat",
    "Mobile web for students (mobile-first only)",
    "Automated ML content moderation",
    "Item insurance or liability features",
]
lines_out = []
for i, s in enumerate(out_scope):
    lines_out.append({"text": f"  •  {s}", "size": 13, "color": SLATE, "space_before": 3 if i > 0 else 2})
add_textbox_rich(sl, lines_out, 7.08, 1.55, 5.6, 5.4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – Features Part 1 (Core)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=SKY)
add_text(sl, "03", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Application Features & Functionalities", 0.5, 0.18, 12, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 3.5, 0.04, fill=AMBER)

add_text(sl, "Core Features", 0.5, 0.88, 12, 0.35, size=14, color=MUTED, italic=True)

cards = [
    ("🔐", "University Login", "Sign in with university email\nJWT session management\nUniversity SSO integration"),
    ("📋", "Report Lost / Found", "Title, category, photos (up to 5)\nCampus map location tagging\nDate, description, tags"),
    ("🔍", "Search & Browse", "Full-text keyword search\nFilters: category, date, location\nPaginated feed sorted by recency"),
    ("🗺️", "Map View", "Items pinned on campus map\nTap pin to see item summary\nFilters applied to map view"),
]

cx, cy, cw, ch = 0.4, 1.35, 3.0, 5.8
gap = 0.2
for i, (icon, title, desc) in enumerate(cards):
    x = cx + i * (cw + gap)
    add_rect(sl, x, cy, cw, ch, fill=WHITE, line_color=BORDER, line_width=Pt(1))
    add_rect(sl, x, cy, cw, 0.06, fill=SKY)  # top accent
    # Icon circle
    add_rect(sl, x + cw/2 - 0.45, cy + 0.25, 0.9, 0.9, fill=BGLT)
    add_text(sl, icon, x + cw/2 - 0.45, cy + 0.28, 0.9, 0.8, size=28, align=PP_ALIGN.CENTER)
    # Title
    add_text(sl, title, x + 0.1, cy + 1.28, cw - 0.2, 0.52, size=14, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    # Divider
    add_rect(sl, x + 0.3, cy + 1.78, cw - 0.6, 0.03, fill=BORDER)
    # Description lines
    for j, dline in enumerate(desc.split("\n")):
        add_text(sl, f"• {dline}", x + 0.18, cy + 1.9 + j*0.58, cw - 0.3, 0.55, size=11.5, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – Features Part 2 (Advanced)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=SKY)
add_text(sl, "03", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Application Features & Functionalities (cont.)", 0.5, 0.18, 12, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 3.5, 0.04, fill=AMBER)

add_text(sl, "Advanced Features", 0.5, 0.88, 12, 0.35, size=14, color=MUTED, italic=True)

cards2 = [
    ("🔔", "Smart Matching", "Auto-match lost & found reports\nCategory + keyword + location scoring\nPush notification within 5 minutes"),
    ("💬", "In-App Messaging", "Secure chat: finder ↔ loser\nPhone/email sharing blocked\nAdmin-reviewable logs"),
    ("✅", "Claim & Verify", "Formal ownership proof required\nFinder approves or rejects claims\nMultiple claims handled fairly"),
    ("🛡️", "Admin Dashboard", "Physical depot management\nContent moderation queue\nPlatform analytics & CSV export"),
]

for i, (icon, title, desc) in enumerate(cards2):
    x = cx + i * (cw + gap)
    add_rect(sl, x, cy, cw, ch, fill=WHITE, line_color=BORDER, line_width=Pt(1))
    add_rect(sl, x, cy, cw, 0.06, fill=AMBER)
    add_rect(sl, x + cw/2 - 0.45, cy + 0.25, 0.9, 0.9, fill=rgb(255,251,235))
    add_text(sl, icon, x + cw/2 - 0.45, cy + 0.28, 0.9, 0.8, size=28, align=PP_ALIGN.CENTER)
    add_text(sl, title, x + 0.1, cy + 1.28, cw - 0.2, 0.52, size=14, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    add_rect(sl, x + 0.3, cy + 1.78, cw - 0.6, 0.03, fill=BORDER)
    for j, dline in enumerate(desc.split("\n")):
        add_text(sl, f"• {dline}", x + 0.18, cy + 1.9 + j*0.58, cw - 0.3, 0.55, size=11.5, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – User Flows
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=SKY)
add_text(sl, "03", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Key User Flows", 0.5, 0.18, 12, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 2.0, 0.04, fill=AMBER)

def draw_flow(sl, title, color, steps, x_start, y_start):
    # Header
    add_rect(sl, x_start, y_start, 12.3, 0.5, fill=color)
    add_text(sl, title, x_start + 0.15, y_start + 0.06, 12.0, 0.38, size=14, bold=True, color=WHITE)

    step_w = 12.3 / len(steps)
    for i, step in enumerate(steps):
        sx = x_start + i * step_w
        sy = y_start + 0.55
        bg = BGLT if color == SKY else rgb(255,251,235)
        add_rect(sl, sx + 0.05, sy, step_w - 0.15, 1.15, fill=bg, line_color=color, line_width=Pt(1))
        add_text(sl, step, sx + 0.1, sy + 0.22, step_w - 0.22, 0.72, size=11, color=SLATE, align=PP_ALIGN.CENTER)
        # Arrow (except last)
        if i < len(steps) - 1:
            ax = sx + step_w - 0.08
            add_text(sl, "→", ax, sy + 0.32, 0.25, 0.5, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)

loser_steps = ["Login", "Report\nLost Item", "Browse\nFound Items", "Submit\nClaim", "Chat with\nFinder", "Confirm\nReceipt ✅"]
finder_steps = ["Login", "Report\nFound Item", "Receive Match\nNotification", "Review\nClaim", "Chat with\nLoser", "Hand Over\nItem ✅"]

draw_flow(sl, "📱  Item Loser Flow", SKY, loser_steps, 0.5, 0.92)
draw_flow(sl, "📱  Item Finder Flow", AMBER, finder_steps, 0.5, 3.35)

# Note
add_rect(sl, 0.5, 5.7, 12.3, 1.55, fill=BGLT)
add_text(sl, "💡  Both flows meet at the Claim & Verification step — ensuring items return to the rightful owner through a structured, privacy-safe process.",
         0.7, 5.82, 12.0, 1.2, size=12.5, italic=True, color=SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – Limitations
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=rgb(239,68,68))
add_text(sl, "04", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Limitations of the Application", 0.5, 0.18, 12, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 2.8, 0.04, fill=rgb(239,68,68))

limits = [
    ("Single campus only", "v1 targets one university; multi-campus support deferred to v2"),
    ("No real-time chat", "Messages use polling (10–30 sec delay); WebSocket deferred to v2"),
    ("No AI image matching", "Matching is rule-based (category + keywords + location); CV deferred"),
    ("Admin adoption risk", "Dashboard value depends on security staff actively using the system"),
    ("Network dependency", "App requires internet connection; offline mode is limited"),
    ("No reward system", "Item returns depend on goodwill; no incentive mechanism in v1"),
    ("Manual moderation", "Chat moderation uses keyword filters + admin review, not ML"),
]

for i, (title, desc) in enumerate(limits):
    row = i % 4
    col = i // 4
    x = 0.5 + col * 6.6
    y = 0.95 + row * 1.6
    w, h = 6.3, 1.45
    add_rect(sl, x, y, w, h, fill=rgb(254,242,242), line_color=rgb(252,165,165), line_width=Pt(1))
    add_rect(sl, x, y, 0.12, h, fill=rgb(239,68,68))
    add_text(sl, f"⚠  {title}", x + 0.25, y + 0.12, w - 0.35, 0.4, size=13, bold=True, color=rgb(185,28,28))
    add_text(sl, desc, x + 0.25, y + 0.55, w - 0.35, 0.75, size=12, color=SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – Technology Stack
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=SKY)
add_text(sl, "05", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Technology Stack", 0.5, 0.18, 12, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 2.2, 0.04, fill=AMBER)

add_text(sl, "All technologies are from the 2190512 Application Development curriculum",
         0.5, 0.88, 12, 0.35, size=12.5, italic=True, color=MUTED)

stack = [
    ("Mobile Frontend",  "React Native (TypeScript)",      "Weeks 12–14", SKY),
    ("Web Frontend",     "HTML + CSS + JavaScript / TS",   "Weeks 10–11", SKY),
    ("Backend API",      "Flask (Python)",                 "Week 7",      rgb(99,102,241)),
    ("Primary Database", "MySQL",                          "Weeks 5–6",   rgb(34,197,94)),
    ("Auth & Real-time", "Firebase (Auth + Firestore)",    "Week 6",      AMBER),
    ("File Storage",     "Firebase Storage",               "Week 6",      AMBER),
    ("Document Store",   "MongoDB (metadata)",             "Week 4",      rgb(34,197,94)),
    ("Containerization", "Docker",                         "Week 7 extra",rgb(99,102,241)),
    ("Admin Analytics",  "Pandas + Plotly",                "Weeks 2–3",   rgb(168,85,247)),
]

cols = [0.5, 5.0, 9.5]
rows_per_col = [3, 3, 3]

for i, (layer, tech, week, color) in enumerate(stack):
    col = i // 3
    row = i % 3
    x = 0.5 + col * 4.25
    y = 1.32 + row * 1.95
    w, h = 3.9, 1.75
    add_rect(sl, x, y, w, h, fill=WHITE, line_color=BORDER, line_width=Pt(1))
    add_rect(sl, x, y, 0.12, h, fill=color)
    # Week chip
    add_rect(sl, x + w - 1.45, y + 0.12, 1.3, 0.3, fill=color)
    add_text(sl, week, x + w - 1.45, y + 0.13, 1.3, 0.28, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, layer, x + 0.25, y + 0.18, w - 1.65, 0.35, size=10.5, bold=False, color=MUTED)
    add_text(sl, tech, x + 0.25, y + 0.55, w - 0.4, 0.5, size=14, bold=True, color=SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – Architecture
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_rect(sl, 0.5, 0.18, 1.0, 0.32, fill=SKY)
add_text(sl, "05", 0.5, 0.2, 1.0, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "System Architecture", 0.5, 0.18, 12, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 2.4, 0.04, fill=AMBER)

# Three horizontal layers
layers = [
    ("Frontend Layer",      [("📱 React Native\n(Mobile App)", SKY),    ("🌐 HTML/CSS/JS\n(Web Admin)", SKY)]),
    ("Backend Layer",       [("⚙️  Flask REST API\n(Python)", rgb(99,102,241))]),
    ("Data Layer",          [("🐬 MySQL\n(Primary DB)", rgb(34,197,94)), ("🍃 MongoDB\n(Metadata)", rgb(34,197,94)), ("🔥 Firebase\n(Auth+Storage+FCM)", AMBER)]),
]

ly_start = [1.0, 2.85, 4.7]
ly_h = 1.6
for li, (lname, boxes) in enumerate(layers):
    ly = ly_start[li]
    # layer bg
    add_rect(sl, 0.3, ly, 12.7, ly_h, fill=rgb(248,250,252), line_color=BORDER, line_width=Pt(1))
    add_text(sl, lname, 0.45, ly + 0.05, 2.2, 0.35, size=10, bold=True, color=MUTED)

    bw = (12.7 - 2.5) / len(boxes)
    for bi, (label, color) in enumerate(boxes):
        bx = 2.8 + bi * (bw + 0.1)
        by = ly + 0.28
        add_rect(sl, bx, by, bw - 0.1, ly_h - 0.45, fill=color)
        add_text(sl, label, bx, by + 0.15, bw - 0.1, ly_h - 0.65, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow to next
    if li < len(layers) - 1:
        add_text(sl, "⬇  REST API calls", 6.0, ly + ly_h + 0.02, 2.0, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Note
add_rect(sl, 0.3, 6.5, 12.7, 0.82, fill=BGLT)
add_text(sl, "📦  Docker containers used for local development & deployment   |   🔑  JWT issued by Flask after Firebase Auth verification   |   📲  FCM push via Firebase",
         0.5, 6.6, 12.4, 0.62, size=11, italic=True, color=SLATE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – Timeline
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, WHITE)
add_rect(sl, 0, 0, 13.33, 0.08, fill=SKY)

add_text(sl, "Project Timeline", 0.5, 0.18, 12, 0.55, size=28, bold=True, color=SLATE, font="Georgia")
add_rect(sl, 0.5, 0.73, 2.4, 0.04, fill=AMBER)

phases = [
    ("Phase 1", "Weeks 1–4", "Setup & Foundation",
     ["University auth (Firebase + JWT)", "Database schema (MySQL + MongoDB)", "Flask API skeleton & Docker setup", "React Native + HTML/JS boilerplate"], SKY),
    ("Phase 2", "Weeks 5–8", "Core Features",
     ["Report lost & found items", "Photo upload (Firebase Storage)", "Search & browse with filters", "Map view (Google Maps SDK)"], rgb(34,197,94)),
    ("Phase 3", "Weeks 9–12", "Advanced Features",
     ["Matching engine (rule-based)", "In-app messaging with moderation", "Claim & verification flow", "Push notifications (FCM)"], rgb(168,85,247)),
    ("Phase 4", "Weeks 13–16", "Polish & Launch",
     ["Admin dashboard & analytics (Plotly)", "UX refinement & testing", "Bug fixes & performance tuning", "Final deployment & presentation"], AMBER),
]

pw = 3.0
gap = 0.1
for i, (phase, weeks, title, tasks, color) in enumerate(phases):
    px = 0.4 + i * (pw + gap)
    py = 0.95

    # Header bar
    add_rect(sl, px, py, pw, 0.72, fill=color)
    add_text(sl, phase, px + 0.1, py + 0.03, pw - 0.2, 0.3, size=11, bold=True, color=WHITE)
    add_text(sl, weeks,  px + 0.1, py + 0.33, pw - 0.2, 0.32, size=10, color=WHITE)

    # Title
    add_rect(sl, px, py + 0.72, pw, 0.55, fill=rgb(248,250,252))
    add_text(sl, title, px + 0.1, py + 0.76, pw - 0.2, 0.44, size=12.5, bold=True, color=SLATE)

    # Tasks
    add_rect(sl, px, py + 1.27, pw, 5.98, fill=WHITE, line_color=BORDER, line_width=Pt(1))
    for j, task in enumerate(tasks):
        ty = py + 1.42 + j * 0.9
        add_rect(sl, px + 0.15, ty, 0.32, 0.32, fill=color)
        add_text(sl, "✓", px + 0.15, ty + 0.02, 0.32, 0.28, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, task, px + 0.6, ty, pw - 0.75, 0.8, size=11.5, color=SLATE)

# Deadline banner
add_rect(sl, 0.4, 6.95, 12.5, 0.4, fill=rgb(254,243,199))
add_text(sl, "📅  Project Proposal Deadline:  26 April 2026 at 23:59", 0.6, 6.97, 12.2, 0.35,
         size=12, bold=True, color=rgb(146,64,14), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – Thank You
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, DARK)

add_rect(sl, 0, 0, 0.45, 7.5, fill=AMBER)
add_rect(sl, 0.45, 3.5, 8, 0.07, fill=AMBER)

add_text(sl, "Thank You", 0.7, 1.2, 9, 1.5, size=64, bold=True, color=WHITE, font="Georgia")
add_text(sl, "UniFind — Reuniting Students with Their Belongings",
         0.7, 2.9, 10, 0.55, size=20, italic=True, color=rgb(186,230,253))

add_text(sl, "Group Members", 0.7, 3.75, 5, 0.38, size=12, bold=True, color=rgb(148,163,184))
add_text(sl, "[Member 1]  |  [Member 2]  |  [Member 3]",
         0.7, 4.15, 8, 0.45, size=16, color=WHITE)

add_text(sl, "2190512 Application Development  |  April 2026",
         0.7, 5.0, 8, 0.4, size=13, color=rgb(100,116,139))

# Decorative right
add_rect(sl, 10.5, 0.5, 2.3, 2.3, fill=rgb(7,89,133))
add_rect(sl, 11.5, 2.5, 1.5, 1.5, fill=SKY)
add_text(sl, "🔍", 11.55, 2.68, 1.4, 1.1, size=42, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
output = "/Users/jirathluekittisup/Uni Stuff/coding/appdev/docs/UniFind_Proposal.pptx"
prs.save(output)
print(f"Saved: {output}")
